"""
Generate a competition presentation deck for the speech disability project.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)


def add_metrics_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Competition Metrics to Present"

    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(12.0)
    height = Inches(4.8)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()

    lines = [
        "Latency: total_ms and stage_timings from /transcribe response",
        "Quality: WER/CER before and after correction",
        "Reliability: confidence score + correction strategy mode",
        "AMD Story: gpu_backend, device, fp16_enabled from runtime telemetry",
        "User Impact: personalization stats and repeat-session improvement",
    ]

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)


def style_deck(prs: Presentation):
    for slide in prs.slides:
        if slide.shapes.title is not None:
            title_tf = slide.shapes.title.text_frame
            for paragraph in title_tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(40)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(22, 43, 98)


def main():
    base_dir = Path(__file__).resolve().parents[1]
    output_dir = base_dir / "demo"
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    add_title_slide(
        prs,
        "AMD-Accelerated Speech Disability Interpreter",
        "Competition Demo Deck | June 2026",
    )

    add_bullet_slide(
        prs,
        "Problem and Vision",
        [
            "Standard ASR fails for dysarthric speech in real scenarios.",
            "Goal: convert impaired speech into reliable text and clear audio.",
            "Differentiator: fast, explainable, and personalized assistive AI.",
        ],
    )

    add_bullet_slide(
        prs,
        "System Architecture",
        [
            "Audio preprocessing -> Whisper ASR -> adaptive correction -> TTS.",
            "FastAPI endpoints for real-time and batch workflows.",
            "Personalization layer learns user-specific mappings over time.",
        ],
    )

    add_bullet_slide(
        prs,
        "What We Implemented for Competition",
        [
            "GPU-aware runtime profiling (device, backend, fp16 capability).",
            "Adaptive correction modes: minimal, dysarthria_sensitive, high_recovery.",
            "Per-stage latency telemetry in API response for benchmarking.",
        ],
    )

    add_bullet_slide(
        prs,
        "Why AMD Infrastructure Matters",
        [
            "ROCm-compatible runtime detection with clean CPU fallback.",
            "Better throughput headroom for batch and streaming workloads.",
            "Clear infrastructure story for judges: performance + accessibility.",
        ],
    )

    add_metrics_slide(prs)

    add_bullet_slide(
        prs,
        "Demo Walkthrough",
        [
            "Upload sample dysarthric audio through /transcribe.",
            "Show original vs corrected text and selected correction strategy.",
            "Play generated corrected audio and highlight latency split.",
        ],
    )

    add_bullet_slide(
        prs,
        "Roadmap",
        [
            "ONNX Runtime ROCm path for correction model acceleration.",
            "Streaming transcription and confidence-aware clarification prompts.",
            "Clinical-grade benchmarking with real dysarthria datasets.",
        ],
    )

    style_deck(prs)

    output_file = output_dir / "competition_presentation.pptx"
    prs.save(str(output_file))
    print(f"Presentation created: {output_file}")


if __name__ == "__main__":
    main()
